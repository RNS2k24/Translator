from PIL import Image
import pytesseract
import pdfplumber
import docx
from pptx import Presentation
import io
import easyocr
import numpy as np

def extract_text(file):
    """
    Extract text from uploaded files including images, PDFs, DOCX, PPTX, and TXT files.
    
    Parameters:
        file: Uploaded file (io.BytesIO) to extract text from.

    Returns:
        str: Extracted text.
    """
    reader = easyocr.Reader(['en'], gpu=False)  # Initialize EasyOCR with English language support

    # Determine file type from filename or file metadata if available
    file_type = getattr(file, "type", None) or getattr(file, "mimetype", None)

    if file_type == "application/pdf":
        text = ""
        with pdfplumber.open(file) as pdf:
            for page in pdf.pages:
                # Extract text from the page if available
                text += page.extract_text() or ""

                # Convert the page to a PIL image
                page_image = page.to_image()
                pil_image = page_image.original  # Get the original PIL image

                # Extract images from the page
                images = page.images
                
                # If images exist on the page, process them
                if images:
                    for img in images:
                        # Get the image's bounding box coordinates
                        x0, top, x1, bottom = img["x0"], img["top"], img["x1"], img["bottom"]
                        
                        # Crop the image out of the page image
                        cropped_image = pil_image.crop((x0, top, x1, bottom))
                        cropped_image_np = np.array(cropped_image)

                        # Extract text from the cropped image using EasyOCR
                        result = reader.readtext(cropped_image_np, detail=0)
                        text += ' '.join(result) + "\n"
        
        return text.strip()

    elif file_type in ["image/png", "image/jpeg", "image/jpg"]:
        # Extract text using EasyOCR for image files
        image = Image.open(file)
        image_np = np.array(image)  # Convert to numpy array for EasyOCR
        result = reader.readtext(image_np, detail=0)
        return ' '.join(result)

    elif file_type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
        doc = docx.Document(file)
        return "\n".join([para.text for para in doc.paragraphs])

    elif file_type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
        text = ""
        prs = Presentation(file)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text

    elif file_type == "text/plain":
        return file.read().decode("utf-8")

    else:
        return "Unsupported file type."
